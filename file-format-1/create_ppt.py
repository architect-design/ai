from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor  # FIXED: Changed RgbColor to RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Wells Fargo Theme Colors
WF_NAVY = RGBColor(0, 48, 87)
WF_GOLD = RGBColor(217, 165, 24)
WF_WHITE = RGBColor(255, 255, 255)
WF_DARK_GRAY = RGBColor(51, 51, 51)
WF_LIGHT_GRAY = RGBColor(240, 240, 240)


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Navy background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WF_NAVY
    shape.line.fill.background()

    # Gold accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.8), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = WF_GOLD
    line.line.fill.background()

    # Title Text
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = WF_WHITE
    p.font.bold = True

    # Subtitle Text
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = WF_GOLD


def add_content_slide(prs, title, bullets, sub_bullets=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Top Navy Header Bar
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = WF_NAVY
    header_bar.line.fill.background()

    # Gold underline for header
    gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.2), Inches(10), Inches(0.05))
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = WF_GOLD
    gold_line.line.fill.background()

    # Title Text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = WF_WHITE
    p.font.bold = True

    # Bullet Points
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = WF_DARK_GRAY
        p.space_after = Pt(12)
        p.level = 0

        # Add sub-bullets if provided
        if sub_bullets and i in sub_bullets:
            for sub in sub_bullets[i]:
                sp = tf2.add_paragraph()
                sp.text = sub
                sp.font.size = Pt(16)
                sp.font.color.rgb = RGBColor(100, 100, 100)  # FIXED
                sp.space_after = Pt(6)
                sp.level = 1
                sp.prBulletFmt.inheritLevelBullet = True


def add_architecture_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = WF_NAVY
    header_bar.line.fill.background()

    gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.2), Inches(10), Inches(0.05))
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = WF_GOLD
    gold_line.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture Diagram"
    p.font.size = Pt(32)
    p.font.color.rgb = WF_WHITE
    p.font.bold = True

    # Architecture Flow as styled text
    arch_text = """[ Web UI Layer ]  (Wells Fargo Themed / Drag & Drop)
       ↓
[ FastAPI Gateway ]  (REST API / File Uploads / Downloads)
       ↓
[ Core Processing Engine ]
    ├── 1. Spec Learner (Regex / Template Parser)
    ├── 2. Validator (Positional Logic / Checksums)
    ├── 3. Generator (Synthetic Data Builder)
    └── 4. DB Transformer (SQL Mapper)
       ↓
[ Knowledge Base ]  (In-Memory Format Rules: NACHA, VISA, GL)
       ↓
[ Database Connectors ]  (Oracle / Postgres / MySQL / Custom)"""

    # Draw a light gray background box for the architecture
    bg_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.4))
    bg_box.fill.solid()
    bg_box.fill.fore_color.rgb = RGBColor(245, 245, 245)  # FIXED
    bg_box.line.color.rgb = WF_GOLD
    bg_box.line.width = Pt(2)

    # Add text
    txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(7.6), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = arch_text
    p2.font.size = Pt(18)
    p2.font.color.rgb = WF_NAVY
    p2.font.name = "Consolas"


def add_demo_steps_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = WF_NAVY
    header_bar.line.fill.background()

    gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.2), Inches(10), Inches(0.05))
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = WF_GOLD
    gold_line.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Live Demo: Solution Steps"
    p.font.size = Pt(32)
    p.font.color.rgb = WF_WHITE
    p.font.bold = True

    steps = [
        ("Step 1: Learn", "Load the built-in NACHA PPD specification into the system memory."),
        ("Step 2: Generate", "Configure batches/entries and generate a compliant test NACHA file."),
        ("Step 3: Validate", "Upload the generated file back into the system to prove 100% validation."),
        ("Step 4: Transform", "Show how database SQL outputs can be mapped to the NACHA format dynamically.")
    ]

    y_pos = 1.6
    for step_title, step_desc in steps:
        # Step Number Box (Gold)
        step_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(2.2),
                                          Inches(0.6))
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = WF_GOLD
        step_box.line.fill.background()
        stf = step_box.text_frame
        stf.paragraphs[0].alignment = PP_ALIGN.CENTER
        stf.paragraphs[0].text = step_title
        stf.paragraphs[0].font.size = Pt(18)
        stf.paragraphs[0].font.bold = True
        stf.paragraphs[0].font.color.rgb = WF_NAVY
        stf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Step Description
        desc_box = slide.shapes.add_textbox(Inches(3.2), Inches(y_pos), Inches(6), Inches(0.6))
        dtf = desc_box.text_frame
        dtf.paragraphs[0].text = step_desc
        dtf.paragraphs[0].font.size = Pt(18)
        dtf.paragraphs[0].font.color.rgb = WF_DARK_GRAY
        dtf.vertical_anchor = MSO_ANCHOR.MIDDLE

        y_pos += 1.2


# --- MAIN EXECUTION ---
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
add_title_slide(prs,
                "Custom Format Intelligence System (CFIS)",
                "Automated File Validation, Generation & Transformation")

# Slide 2: Problem Statement
add_content_slide(prs, "The Problem Statement", [
    "Manual Validation is Error-Prone: Testing banking files (NACHA, VISA VCF) manually leads to high error rates in production.",
    "Rigid Formatting Rules: Positional file formats require exact character counts, routing checksums, and cross-record totals.",
    "Database-to-File Friction: Extracting data from Oracle/Postgres and mapping it to strict flat-file layouts is slow and requires custom scripts.",
    "Lack of Dynamic Learning: Existing tools require hard-coded logic for every new file specification introduced by clients."
])

# Slide 3: Our Solution
add_content_slide(prs, "The CFIS Solution", [
    "Rule-Based 'Learning' Engine: Parses specification documents and builds internal rule trees dynamically—no neural networks, 100% deterministic.",
    "Multi-Format Support: Out-of-the-box support for NACHA, VISA VCF, Oracle GL, and custom formats.",
    "Automated Test Data Generation: Creates mathematically valid synthetic data (e.g., valid ABA routing checksums).",
    "Seamless DB Transformation: Connect to any enterprise DB, run SQL, and instantly convert result sets into perfectly formatted flat files."
])

# Slide 4: Architecture
add_architecture_slide(prs)

# Slide 5: Design & Capabilities
add_content_slide(prs, "System Design & Capabilities", [
    "Specification Learner UI: Drag-and-drop PDFs/Excel specs or load built-in templates.",
    "Real-Time Validator: Upload a file and instantly get a score (e.g., 98.5%), with exact error locations (Record 45, Position 2-3).",
    "Test Data Factory: Generate edge-case testing files (max amounts, special characters, boundary lengths).",
    "Database Integration: Securely store DB connections, write SQL, map columns to file positions visually."
])

# Slide 6: Demo Steps
add_demo_steps_slide(prs)

# Slide 7: Conclusion
add_title_slide(prs,
                "Benefits & Conclusion",
                "Secure • Deterministic • Easily Extensible • Zero External API Dependencies")

# Save
prs.save('CFIS_Demo_Presentation.pptx')
print("Presentation saved successfully as 'CFIS_Demo_Presentation.pptx'")